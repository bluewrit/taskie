"""File preview engine.

Classifies an uploaded file and produces a preview payload for the frontend:
  image / pdf / video / audio  -> streamed inline by the browser
  text / code / markdown       -> raw text content
  csv / tsv                    -> parsed table (columns + rows)
  json                         -> pretty-printed JSON
  docx / pptx / pdf            -> extracted text ("document" preview)
  xlsx                         -> sheets parsed into tables
  zip                          -> archive entry listing
  anything else                -> download fallback
"""
import csv
import io
import json
import zipfile
from pathlib import Path

MAX_TEXT_BYTES = 200_000
MAX_DOC_CHARS = 60_000

TEXT_EXTENSIONS = {
    "txt", "md", "markdown", "log", "rst", "ini", "cfg", "conf", "yaml", "yml",
    "toml", "env", "properties", "xml", "html", "htm", "css", "scss", "less",
    "js", "jsx", "ts", "tsx", "vue", "svelte", "py", "rb", "go", "rs", "java",
    "kt", "c", "h", "cpp", "hpp", "cs", "php", "swift", "sql", "sh", "bash",
    "zsh", "ps1", "bat", "csv", "tsv", "json", "jsonl", "ndjson", "ipynb",
    "tex", "r", "lua", "pl", "dart", "gradle", "dockerfile", "makefile", "gitignore",
}

LANGUAGE_BY_EXT = {
    "py": "python", "js": "javascript", "jsx": "jsx", "ts": "typescript",
    "tsx": "tsx", "html": "html", "css": "css", "json": "json", "md": "markdown",
    "sql": "sql", "sh": "bash", "yaml": "yaml", "yml": "yaml", "java": "java",
    "go": "go", "rs": "rust", "rb": "ruby", "php": "php", "c": "c", "cpp": "cpp",
    "cs": "csharp", "xml": "xml", "txt": "plaintext", "log": "plaintext",
}

IMAGE_EXTS = {"png", "jpg", "jpeg", "gif", "webp", "bmp", "svg", "ico", "avif"}
VIDEO_EXTS = {"mp4", "webm", "ogv", "mov"}
AUDIO_EXTS = {"mp3", "wav", "ogg", "m4a", "aac", "flac", "opus"}
ARCHIVE_EXTS = {"zip"}


def preview_kind(filename: str, mime_type: str) -> str:
    """Return the preview family for a file."""
    ext = Path(filename).suffix.lower().lstrip(".")
    if ext in IMAGE_EXTS or (mime_type or "").startswith("image/"):
        return "image"
    if ext == "pdf" or mime_type == "application/pdf":
        return "pdf"
    if ext in VIDEO_EXTS or (mime_type or "").startswith("video/"):
        return "video"
    if ext in AUDIO_EXTS or (mime_type or "").startswith("audio/"):
        return "audio"
    if ext in ("docx", "pptx"):
        return "document"
    if ext == "xlsx":
        return "sheets"
    if ext in ("csv", "tsv"):
        return "table"
    if ext in ARCHIVE_EXTS:
        return "archive"
    if ext in TEXT_EXTENSIONS or ext == "" or (mime_type or "").startswith("text/"):
        return "text"
    return "download"


def _truncate(text: str) -> str:
    if len(text) > MAX_DOC_CHARS:
        return text[:MAX_DOC_CHARS] + f"\n\n… truncated ({len(text) - MAX_DOC_CHARS} characters omitted)"
    return text


def _extract_docx(path: Path) -> str:
    import docx  # python-docx

    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        parts.append("")
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return _truncate("\n".join(parts) or "(empty document)")


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"── Slide {i} ──"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)
        blocks.append("\n".join(lines))
    return _truncate("\n\n".join(blocks) or "(empty presentation)")


def _extract_pdf_text(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages[:40], 1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append(f"── Page {i} ──\n{text}")
    return _truncate("\n\n".join(pages) or "(no extractable text — use the PDF viewer)")


def _extract_xlsx(path: Path) -> list[dict]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            if len(rows) >= 200:
                break
            rows.append(["" if v is None else str(v) for v in row])
        if not rows:
            sheets.append({"name": ws.title, "columns": [], "rows": []})
            continue
        columns = rows[0]
        sheets.append({"name": ws.title, "columns": columns, "rows": rows[1:]})
    wb.close()
    return sheets


def _extract_zip(path: Path) -> list[dict]:
    with zipfile.ZipFile(str(path)) as zf:
        entries = []
        for info in zf.infolist()[:500]:
            entries.append({
                "name": info.filename,
                "size": info.file_size,
                "is_dir": info.is_dir(),
            })
        return entries


def _parse_csv(path: Path, delimiter: str) -> dict:
    with open(path, "r", encoding="utf-8", errors="replace", newline="") as fh:
        sample = fh.read(4096)
        fh.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
        except csv.Error:
            delimiter = delimiter
        reader = csv.reader(fh, delimiter=delimiter)
        rows = []
        for i, row in enumerate(reader):
            if i >= 200:
                break
            rows.append(row)
    if not rows:
        return {"columns": [], "rows": []}
    return {"columns": rows[0], "rows": rows[1:]}


def build_preview(path: Path, filename: str, mime_type: str) -> dict:
    """Return a JSON-serialisable preview payload for a stored file."""
    kind = preview_kind(filename, mime_type)
    ext = Path(filename).suffix.lower().lstrip(".")
    payload: dict = {"kind": kind, "filename": filename}

    if kind in ("image", "pdf", "video", "audio", "download"):
        # Browser handles these via the raw stream endpoint; nothing to extract.
        payload["extracted_text"] = None
        if kind == "pdf":
            try:
                payload["extracted_text"] = _extract_pdf_text(path)
            except Exception as exc:  # pragma: no cover - defensive
                payload["extracted_text"] = f"(text extraction failed: {exc})"
        return payload

    if kind == "document":
        try:
            if ext == "docx":
                payload["content"] = _extract_docx(path)
            else:
                payload["content"] = _extract_pptx(path)
        except Exception as exc:
            payload["content"] = f"(could not extract document text: {exc})"
        return payload

    if kind == "sheets":
        try:
            payload["sheets"] = _extract_xlsx(path)
        except Exception as exc:
            payload["sheets"] = []
            payload["error"] = str(exc)
        return payload

    if kind == "table":
        try:
            table = _parse_csv(path, "," if ext == "csv" else "\t")
            payload.update(table)
        except Exception as exc:
            payload["columns"], payload["rows"] = [], []
            payload["error"] = str(exc)
        return payload

    if kind == "archive":
        try:
            payload["entries"] = _extract_zip(path)
        except Exception as exc:
            payload["entries"] = []
            payload["error"] = str(exc)
        return payload

    # text family
    raw = path.read_bytes()[:MAX_TEXT_BYTES]
    if ext == "json":
        try:
            raw = json.dumps(json.loads(raw.decode("utf-8")), indent=2).encode("utf-8")
        except Exception:
            pass
    text = raw.decode("utf-8", errors="replace")
    payload["content"] = text
    payload["language"] = LANGUAGE_BY_EXT.get(ext, "plaintext")
    if Path(filename).name.lower() == "dockerfile":
        payload["language"] = "dockerfile"
    return payload


def is_binary_likely(path: Path) -> bool:
    chunk = path.open("rb").read(1024)
    return b"\x00" in chunk or not chunk
